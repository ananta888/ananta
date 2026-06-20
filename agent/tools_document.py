from flask import current_app

from agent.tools_registry import registry


@registry.register(
    name="doc_extract",
    description="Extrahiert Text aus einem Dokument (PDF, DOCX, XLSX, PPTX).",
    parameters={
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Pfad zum Dokument"},
        },
        "required": ["path"],
    },
)
def doc_extract_tool(path: str):
    from agent.services.platform_governance_service import get_platform_governance_service
    gov = get_platform_governance_service()
    cfg = current_app.config.get("AGENT_CONFIG")

    if not gov.evaluate_action_pack_access("document", cfg):
        return {"error": "Action Pack 'document' ist deaktiviert."}

    import os
    if not os.path.exists(path):
        return {"error": f"Datei '{path}' nicht gefunden."}

    ext = os.path.splitext(path)[1].lower()
    text = ""

    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(path)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted: text += extracted + "\n"
        elif ext == ".docx":
            import docx
            doc = docx.Document(path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
        elif ext == ".pptx":
            import pptx
            prs = pptx.Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            return {"error": f"Unterstuetztes Format für '{ext}' nicht gefunden."}

        from agent.common.audit import log_audit
        log_audit("doc_extract", {"path": path, "size": len(text)})

        return {
            "content": text[:30000],
            "path": path,
            "format": ext.lstrip(".")
        }
    except Exception as e:
        return {"error": f"Fehler bei Extraktion aus {path}: {e}"}
