from __future__ import annotations

import json
import logging
import re
import zipfile
from pathlib import Path

try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover - optional dependency
    DocxDocument = None

try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover - optional dependency
    load_workbook = None

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover - optional dependency
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None

try:
    from lxml import html
except ImportError:  # pragma: no cover - optional dependency
    html = None


class ExtractionService:
    """Extracts text when feasible and falls back to metadata-only modes."""

    FULLY_INDEXED_EXTENSIONS = {
        ".md", ".rst", ".json", ".yaml", ".yml", ".xml", ".csv",
        ".py", ".ts", ".tsx", ".js", ".jsx", ".java", ".go", ".rs", ".c", ".cpp", ".cs", ".php",
    }
    TEXT_EXTRACTED_EXTENSIONS = {".txt", ".log", ".ini", ".toml", ".html", ".htm"}
    OFFICE_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp"}
    BINARY_REFERENCE_EXTENSIONS = {
        ".png", ".jpg", ".jpeg", ".webp",
        ".mp3", ".wav", ".mp4", ".mkv",
        ".zip", ".tar", ".gz", ".eml", ".msg",
    }

    def _base_metadata(self, *, storage_path: str, filename: str, media_type: str, suffix: str) -> dict:
        return {
            "filename": filename,
            "media_type": media_type,
            "storage_path": storage_path,
            "suffix": suffix,
        }

    def _read_text(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="ignore")

    def _html_text(self, path: Path) -> str:
        if html is None:
            return "lxml_unavailable"
        raw = self._read_text(path)
        document = html.fromstring(raw or "<html></html>")
        for node in document.xpath("//script|//style"):
            node.drop_tree()
        text = " ".join(part.strip() for part in document.xpath("//text()") if str(part).strip())
        return re.sub(r"\s+", " ", text).strip()

    def _pdf_text(self, path: Path) -> tuple[str | None, str | None]:
        if PdfReader is None:
            return None, "pdf_extractor_unavailable"
        try:
            reader = PdfReader(str(path))
            text = "\n".join((page.extract_text() or "").strip() for page in reader.pages).strip()
            return text or None, "pdf_text_extracted" if text else "pdf_no_text"
        except Exception as exc:  # pragma: no cover - parser behavior varies by fixture
            logging.warning("PDF extraction failed for %s: %s", path, exc)
            return None, "pdf_extraction_failed"

    def _docx_text(self, path: Path) -> tuple[str | None, str | None]:
        if DocxDocument is None:
            return None, "docx_extractor_unavailable"
        try:
            document = DocxDocument(str(path))
            text = "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()).strip()
            return text or None, "docx_text_extracted" if text else "docx_no_text"
        except Exception as exc:  # pragma: no cover - parser behavior varies by fixture
            logging.warning("DOCX extraction failed for %s: %s", path, exc)
            return None, "docx_extraction_failed"

    def _xlsx_text(self, path: Path) -> tuple[str | None, str | None]:
        if load_workbook is None:
            return None, "xlsx_extractor_unavailable"
        try:
            workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
            lines: list[str] = []
            for sheet in workbook.worksheets:
                lines.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(cell).strip() for cell in row if cell not in (None, "")]
                    if values:
                        lines.append(" | ".join(values))
            text = "\n".join(lines).strip()
            return text or None, "xlsx_text_extracted" if text else "xlsx_no_text"
        except Exception as exc:  # pragma: no cover - parser behavior varies by fixture
            logging.warning("XLSX extraction failed for %s: %s", path, exc)
            return None, "xlsx_extraction_failed"

    def _pptx_text(self, path: Path) -> tuple[str | None, str | None]:
        if Presentation is None:
            return None, "pptx_extractor_unavailable"
        try:
            presentation = Presentation(str(path))
            lines: list[str] = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                lines.append(f"# Slide {slide_index}")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and str(text).strip():
                        lines.append(str(text).strip())
            text = "\n".join(lines).strip()
            return text or None, "pptx_text_extracted" if text else "pptx_no_text"
        except Exception as exc:  # pragma: no cover - parser behavior varies by fixture
            logging.warning("PPTX extraction failed for %s: %s", path, exc)
            return None, "pptx_extraction_failed"

    def _opendocument_text(self, path: Path) -> tuple[str | None, str | None]:
        try:
            with zipfile.ZipFile(path, "r") as archive:
                raw = archive.read("content.xml").decode("utf-8", errors="ignore")
        except Exception as exc:  # pragma: no cover - parser behavior varies by fixture
            logging.warning("OpenDocument extraction failed for %s: %s", path, exc)
            return None, "opendocument_extraction_failed"
        text = re.sub(r"<[^>]+>", " ", raw)
        text = re.sub(r"\s+", " ", text).strip()
        return text or None, "opendocument_text_extracted" if text else "opendocument_no_text"

    def _extract_office_text(self, path: Path, suffix: str) -> tuple[str | None, str]:
        if suffix == ".pdf":
            return self._pdf_text(path)
        if suffix == ".docx":
            return self._docx_text(path)
        if suffix == ".xlsx":
            return self._xlsx_text(path)
        if suffix == ".pptx":
            return self._pptx_text(path)
        if suffix in {".odt", ".ods", ".odp"}:
            return self._opendocument_text(path)
        return None, "office_extraction_not_enabled"

    def extract(self, *, storage_path: str, filename: str, media_type: str) -> dict:
        path = Path(storage_path)
        suffix = path.suffix.lower() or Path(filename).suffix.lower()
        metadata = self._base_metadata(storage_path=storage_path, filename=filename, media_type=media_type, suffix=suffix)

        if (
            suffix in self.FULLY_INDEXED_EXTENSIONS
            or media_type in {"application/json", "application/xml"}
            or media_type.startswith("application/x-python")
        ):
            text_content = self._read_text(path)
            mode = "fully-indexed"
            if suffix == ".json":
                try:
                    parsed = json.loads(text_content)
                    metadata["json_root_type"] = type(parsed).__name__
                except Exception:
                    metadata["json_root_type"] = "invalid"
            metadata["content_family"] = "structured_text"
            return {
                "extraction_status": "completed",
                "extraction_mode": mode,
                "text_content": text_content,
                "metadata": metadata,
            }

        if suffix in {".html", ".htm"} or media_type.startswith("text/html"):
            text_content = self._html_text(path)
            return {
                "extraction_status": "completed",
                "extraction_mode": "text-extracted",
                "text_content": text_content or None,
                "metadata": {**metadata, "content_family": "html_document", "reason": "html_text_extracted"},
            }

        if suffix in self.TEXT_EXTRACTED_EXTENSIONS or media_type.startswith("text/"):
            text_content = self._read_text(path)
            return {
                "extraction_status": "completed",
                "extraction_mode": "text-extracted",
                "text_content": text_content,
                "metadata": {**metadata, "content_family": "plain_text"},
            }

        if suffix in self.OFFICE_EXTENSIONS:
            text_content, reason = self._extract_office_text(path, suffix)
            if text_content:
                return {
                    "extraction_status": "completed",
                    "extraction_mode": "text-extracted",
                    "text_content": text_content,
                    "metadata": {**metadata, "content_family": "office_document", "reason": reason},
                }
            return {
                "extraction_status": "completed",
                "extraction_mode": "metadata-only",
                "text_content": None,
                "metadata": {**metadata, "content_family": "office_document", "reason": reason},
            }

        if suffix in self.BINARY_REFERENCE_EXTENSIONS:
            return {
                "extraction_status": "completed",
                "extraction_mode": "raw-only",
                "text_content": None,
                "metadata": {**metadata, "content_family": "binary_reference", "reason": "binary_or_reference_only"},
            }

        return {
            "extraction_status": "completed",
            "extraction_mode": "raw-only",
            "text_content": None,
            "metadata": {**metadata, "content_family": "unknown_binary", "reason": "binary_or_unsupported_format"},
        }


extraction_service = ExtractionService()


def get_extraction_service() -> ExtractionService:
    return extraction_service
